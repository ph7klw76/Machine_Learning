
import openai

openai.api_key = "sk-xxxxxxxxxxx"
prompt = "create a lecture note on Introduction to Superconductors for undergraduate physics final year students"
response = openai.Completion.create(
  model="text-davinci-003",
  prompt=prompt,
  temperature=0.75,
  max_tokens=2060,
  top_p=0.75,
  frequency_penalty=0.0,
  presence_penalty=0.0
)
response=response['choices'][0]['text']
response=response.split('\n\n')
no=len(response)
from pptx import Presentation
prs = Presentation() 
title_slide_layout = prs.slide_layouts[1]
for i in range(1,no-1):
    slide=prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1] #placeholder for subtitle
    title.text = response[1]
    subtitle.text = response[i+1] # subtitle

prs.save('E:/example2.pptx')
